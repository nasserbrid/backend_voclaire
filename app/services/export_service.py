import io
from typing import Optional

from docx import Document
from fpdf import FPDF
from pptx import Presentation

from app.logger import logger


def generate_docx_free(
    transcription_text: str,
    file_name: str,
    segments: Optional[list[dict]] = None,
) -> bytes:
    doc = Document()
    doc.add_heading(file_name, level=1)

    if segments:
        speaker_order: dict[str, int] = {}
        for segment in segments:
            speaker = segment["speaker"]
            if speaker not in speaker_order:
                speaker_order[speaker] = len(speaker_order)

        for segment in segments:
            speaker_index = speaker_order[segment["speaker"]]
            minutes = int(segment["start"] // 60)
            seconds = int(segment["start"] % 60)
            timestamp = f"{minutes}:{seconds:02d}"

            speaker_paragraph = doc.add_paragraph()
            speaker_run = speaker_paragraph.add_run(f"Intervenant {speaker_index + 1} · {timestamp}")
            speaker_run.bold = True

            doc.add_paragraph(segment["text"])
    else:
        doc.add_paragraph(transcription_text)

    buffer = io.BytesIO()
    doc.save(buffer)
    file_bytes = buffer.getvalue()
    logger.info(f"DOCX free généré : {len(file_bytes)} bytes")
    return file_bytes


def generate_docx_pro(structured: dict, file_name: str) -> bytes:
    doc = Document()
    doc.add_heading(structured["titre"], level=1)
    doc.add_heading("Introduction", level=2)
    doc.add_paragraph(structured["introduction"])
    doc.add_heading("Points abordés", level=2)
    points: list[str] = structured["points"]
    for index, point in enumerate(points, start=1):
        doc.add_paragraph(f"{index}. {point}")
    doc.add_heading("Conclusion", level=2)
    doc.add_paragraph(structured["conclusion"])
    buffer = io.BytesIO()
    doc.save(buffer)
    file_bytes = buffer.getvalue()
    logger.info(f"DOCX pro généré : {len(file_bytes)} bytes")
    return file_bytes


def generate_pdf_pro(structured: dict, file_name: str) -> bytes:
    pdf = FPDF()
    pdf.add_page()

    pdf.set_font("Helvetica", style="B", size=16)
    pdf.multi_cell(0, 10, structured["titre"])
    pdf.ln(4)

    pdf.set_font("Helvetica", style="B", size=13)
    pdf.multi_cell(0, 8, "Introduction")
    pdf.set_font("Helvetica", size=11)
    pdf.multi_cell(0, 7, structured["introduction"])
    pdf.ln(4)

    pdf.set_font("Helvetica", style="B", size=13)
    pdf.multi_cell(0, 8, "Points abordés")
    pdf.set_font("Helvetica", size=11)
    points: list[str] = structured["points"]
    for index, point in enumerate(points, start=1):
        pdf.multi_cell(0, 7, f"{index}. {point}")
    pdf.ln(4)

    pdf.set_font("Helvetica", style="B", size=13)
    pdf.multi_cell(0, 8, "Conclusion")
    pdf.set_font("Helvetica", size=11)
    pdf.multi_cell(0, 7, structured["conclusion"])

    file_bytes = bytes(pdf.output())
    logger.info(f"PDF pro généré : {len(file_bytes)} bytes")
    return file_bytes


def generate_pptx_pro(structured: dict, file_name: str) -> bytes:
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text_frame.text = structured["titre"]
    title_slide.placeholders[1].text_frame.text = file_name

    content_slide_layout = prs.slide_layouts[1]

    intro_slide = prs.slides.add_slide(content_slide_layout)
    intro_slide.shapes.title.text_frame.text = "Introduction"
    intro_slide.placeholders[1].text_frame.text = structured["introduction"]

    points: list[str] = structured["points"]
    for index, point in enumerate(points, start=1):
        point_slide = prs.slides.add_slide(content_slide_layout)
        point_slide.shapes.title.text_frame.text = f"Point {index}"
        point_slide.placeholders[1].text_frame.text = point

    conclusion_slide = prs.slides.add_slide(content_slide_layout)
    conclusion_slide.shapes.title.text_frame.text = "Conclusion"
    conclusion_slide.placeholders[1].text_frame.text = structured["conclusion"]

    buffer = io.BytesIO()
    prs.save(buffer)
    file_bytes = buffer.getvalue()
    logger.info(f"PPTX pro généré : {len(file_bytes)} bytes")
    return file_bytes
